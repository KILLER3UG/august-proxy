"""Artifact + circuit tool smoke tests (offline).

``create_pptx`` / ``render_chart`` / ``render_video`` / ``draw_circuit``
write real files; each test verifies the artifact is a valid file of the
expected kind. ``simulate_circuit`` is exercised for its ngspice-missing
guidance path (and the happy path when ngspice happens to be installed).
``search_component`` is network-bound — only the input-validation path is
tested here.
"""

from __future__ import annotations

import pytest
from app.services.tools import artifact_tools, circuit_tools


def test_create_pptx_writes_deck(tmp_path):
    out = tmp_path / 'deck.pptx'
    result = artifact_tools.create_pptx(
        str(out),
        [
            {'title': 'Deck', 'bullets': ['one', 'two'], 'notes': 'note'},
            'Section title',
        ],
        workspace=str(tmp_path),
    )
    assert result['slideCount'] == 2
    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 2
    assert prs.slides[0].shapes.title.text == 'Deck'


def test_create_pptx_rejects_empty_and_bad_ext(tmp_path):
    with pytest.raises(ValueError):
        artifact_tools.create_pptx(str(tmp_path / 'deck.pptx'), [], workspace=str(tmp_path))
    with pytest.raises(ValueError):
        artifact_tools.create_pptx(str(tmp_path / 'deck.txt'), [{'title': 'x'}], workspace=str(tmp_path))


def test_create_html_artifact_writes_document(tmp_path):
    out = tmp_path / 'explainer.html'
    doc = (
        '<!doctype html><html><head><style>body{background:#111}</style></head>'
        '<body><canvas id="c"></canvas><script>const c=document.getElementById("c");</script></body></html>'
    )
    result = artifact_tools.create_html_artifact(str(out), doc, workspace=str(tmp_path))
    assert result['bytes'] > 0
    assert out.exists()
    assert 'canvas' in out.read_text(encoding='utf-8')


def test_create_html_artifact_adds_title_when_missing(tmp_path):
    out = tmp_path / 't.html'
    artifact_tools.create_html_artifact(
        str(out), '<html><head></head><body>hi</body></html>',
        title='Embeddings 101', workspace=str(tmp_path),
    )
    text = out.read_text(encoding='utf-8')
    assert '<title>Embeddings 101</title>' in text


def test_create_html_artifact_flags_external_refs(tmp_path):
    out = tmp_path / 'ext.html'
    result = artifact_tools.create_html_artifact(
        str(out),
        '<html><body><script src="https://cdn.example.com/x.js"></script></body></html>',
        workspace=str(tmp_path),
    )
    assert result['externalRefs'], 'remote script should be flagged'


def test_create_html_artifact_rejects_empty_and_bad_ext(tmp_path):
    import pytest as _pytest
    with _pytest.raises(ValueError):
        artifact_tools.create_html_artifact(str(tmp_path / 'a.html'), '', workspace=str(tmp_path))
    with _pytest.raises(ValueError):
        artifact_tools.create_html_artifact(str(tmp_path / 'a.txt'), '<p>x</p>', workspace=str(tmp_path))


def test_render_chart_png(tmp_path):
    out = tmp_path / 'chart.png'
    result = artifact_tools.render_chart(
        str(out), 'line', [[1, 3, 2, 5]], title='t', workspace=str(tmp_path)
    )
    assert result['kind'] == 'line'
    assert out.read_bytes()[:8] == b'\x89PNG\r\n\x1a\n'


def test_render_chart_traces_dict_and_file(tmp_path):
    import json as _json

    traces = {
        'v(out)': {'x': [0.0, 1e-3, 2e-3], 'y': [0.0, 0.63, 0.86], 'xunit': 's', 'unit': 'V'},
        'v(in)': {'x': [0.0, 1e-3, 2e-3], 'y': [5.0, 5.0, 5.0], 'xunit': 's', 'unit': 'V'},
    }
    out = tmp_path / 'waves.png'
    result = artifact_tools.render_chart(
        str(out), 'line', traces=traces, title='RC step', workspace=str(tmp_path)
    )
    assert result['traceNames'] == ['v(out)', 'v(in)']
    assert out.read_bytes()[:8] == b'\x89PNG\r\n\x1a\n'
    # File form: circuit_simulate's tracesFile path round-trips the same.
    tf = tmp_path / 'sim_traces.json'
    tf.write_text(_json.dumps(traces), encoding='utf-8')
    out2 = tmp_path / 'waves2.png'
    result2 = artifact_tools.render_chart(
        str(out2), 'line', traces='sim_traces.json', workspace=str(tmp_path)
    )
    assert result2['traceNames'] == ['v(out)', 'v(in)']


def test_render_chart_traces_validation(tmp_path):
    good = {'a': {'x': [0.0, 1.0], 'y': [0.0, 1.0]}}
    # Traces only plot as lines.
    with pytest.raises(ValueError):
        artifact_tools.render_chart(
            str(tmp_path / 'x.png'), 'bar', traces=good, workspace=str(tmp_path)
        )
    # Mismatched x/y lengths are rejected.
    with pytest.raises(ValueError):
        artifact_tools.render_chart(
            str(tmp_path / 'y.png'), 'line',
            traces={'a': {'x': [0.0, 1.0], 'y': [0.0]}}, workspace=str(tmp_path),
        )
    # Missing traces file is rejected.
    with pytest.raises(ValueError):
        artifact_tools.render_chart(
            str(tmp_path / 'z.png'), 'line', traces='nope.json', workspace=str(tmp_path)
        )
    # Neither series nor traces is rejected.
    with pytest.raises(ValueError):
        artifact_tools.render_chart(
            str(tmp_path / 'w.png'), 'line', workspace=str(tmp_path)
        )


def test_render_video_mp4(tmp_path):
    from PIL import Image

    frames = []
    for i in range(3):
        p = tmp_path / f'f{i}.png'
        Image.new('RGB', (64, 48), (i * 40, 0, 0)).save(p)
        frames.append(str(p))
    out = tmp_path / 'clip.mp4'
    result = artifact_tools.render_video(
        str(out), frames, fps=6, hold_last_ms=100, workspace=str(tmp_path)
    )
    assert result['frameCount'] == 3
    assert out.stat().st_size > 0
    # MP4 signature (ftyp box) in the first bytes.
    assert out.read_bytes()[4:8] == b'ftyp'


def test_draw_circuit_png(tmp_path):
    out = tmp_path / 'sch.png'
    result = artifact_tools.draw_circuit(
        str(out),
        [
            {'type': 'battery', 'label': '9V', 'dir': 'right'},
            {'type': 'resistor', 'label': '1k', 'dir': 'down'},
            {'type': 'ground', 'dir': 'left'},
        ],
        workspace=str(tmp_path),
    )
    assert result['elementCount'] == 3
    assert out.read_bytes()[:8] == b'\x89PNG\r\n\x1a\n'


def test_draw_circuit_unknown_element(tmp_path):
    with pytest.raises(ValueError, match='Unknown element type'):
        artifact_tools.draw_circuit(
            str(tmp_path / 'x.png'), [{'type': 'flux_capacitor'}], workspace=str(tmp_path)
        )


@pytest.mark.asyncio
async def test_simulate_circuit_without_ngspice(monkeypatch, tmp_path):
    monkeypatch.setattr(circuit_tools, 'resolve_ngspice', lambda: None)
    result = await circuit_tools.simulate_circuit('V1 0 0 0\n.end', workspace=str(tmp_path))
    assert result.get('installed') is False
    assert 'ngspice' in result.get('error', '')


@pytest.mark.asyncio
async def test_simulate_circuit_happy_path(tmp_path):
    exe = circuit_tools.resolve_ngspice()
    if exe is None:
        pytest.skip('ngspice not installed on this machine')
    netlist = '\n'.join(
        [
            '* divider',
            'V1 in 0 DC 10',
            'R1 in out 1k',
            'R2 out 0 1k',
            '.control',
            'op',
            'print v(out)',
            '.endc',
            '.end',
        ]
    )
    result = await circuit_tools.simulate_circuit(netlist, name='divider', workspace=str(tmp_path))
    assert result['exitCode'] == 0
    assert result['measures'], f'expected parsed measures, got: {result["logTail"][-400:]}'


def test_search_component_requires_query():
    import asyncio

    result = asyncio.run(circuit_tools.search_component(''))
    assert 'error' in result
