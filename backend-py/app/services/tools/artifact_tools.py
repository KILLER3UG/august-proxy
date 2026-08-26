"""Artifact creation tools — PPTX decks, chart images, videos, circuits.

Complements the audit-side ``pptx_tools`` (inspect/comment): these BUILD
files in the session workspace. Everything is offline:

  * ``create_pptx``    — python-pptx title+bullet decks
  * ``render_chart``   — matplotlib (Agg) → PNG
  * ``render_video``   — imageio + bundled ffmpeg → MP4 (images→video, or
                         frame-matrix animations)
  * ``draw_circuit``   — schemdraw → PNG schematic

All writers bind their output path to the workspace like every other
mutating tool (sandbox.paths.bind_path), and accept either structured
lists or JSON strings (models frequently stringify arguments).
"""

from __future__ import annotations

import json
import logging
from pathlib import Path
from typing import Any

from app.services.sandbox.paths import bind_path

logger = logging.getLogger(__name__)

# Matplotlib must pick the headless backend BEFORE pyplot is imported;
# importing this module lazily inside each call keeps startup light.
_AGG_READY = False


def _ensure_matplotlib() -> Any:
    global _AGG_READY
    if not _AGG_READY:
        import matplotlib

        matplotlib.use('Agg', force=True)
        _AGG_READY = True
    import matplotlib.pyplot as plt

    return plt


def _as_list(value: Any) -> list[Any]:
    """Accept JSON strings or native lists for model-supplied payloads."""
    if isinstance(value, str):
        try:
            value = json.loads(value)
        except json.JSONDecodeError as exc:
            raise ValueError(f'Invalid JSON payload: {exc}') from exc
    if not isinstance(value, list):
        raise ValueError('Expected a list (or a JSON-encoded list).')
    return value


def _bind_write(path: str, workspace: str, suffixes: tuple[str, ...]) -> str:
    """Workspace-bind an output path and enforce the extension."""
    if not (path or '').strip():
        raise ValueError('path is required (relative to the workspace).')
    bound, err = bind_path(path, workspace, for_write=True)
    if err or bound is None:
        raise ValueError(err or f'Invalid path: {path}')
    if bound.suffix.lower() not in suffixes:
        raise ValueError(
            f'{path} must end with one of {" / ".join(suffixes)}'
        )
    bound.parent.mkdir(parents=True, exist_ok=True)
    return str(bound)


# ── PPTX ──────────────────────────────────────────────────────────────────


def _coerce_slide(raw: Any, index: int) -> dict[str, Any]:
    if isinstance(raw, str):
        raw = {'title': raw}
    if not isinstance(raw, dict):
        raise ValueError(f'slides[{index}] must be an object or a title string.')
    return raw


def create_pptx(path: str, slides: Any, workspace: str = '') -> dict[str, Any]:
    """Build a deck from [{'title': str, 'bullets': [str], 'notes': str}]."""
    out = _bind_write(path, workspace, ('.pptx',))
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    items = _as_list(slides)
    if not items:
        raise ValueError('slides must contain at least one slide.')

    made = 0
    for i, raw in enumerate(items):
        spec = _coerce_slide(raw, i)
        layout = prs.slide_layouts[1 if spec.get('bullets') else 0]
        slide = prs.slides.add_slide(layout)
        title_text = str(spec.get('title') or f'Slide {i + 1}')
        slide.shapes.title.text = title_text
        body = slide.placeholders[1] if spec.get('bullets') else None
        if body is not None:
            tf = body.text_frame
            tf.clear()
            bullets = [str(b) for b in _as_list(spec.get('bullets')) if str(b).strip()]
            for j, bullet in enumerate(bullets):
                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.text = bullet
                para.font.size = Pt(18)
                para.level = 0
        notes = str(spec.get('notes') or '').strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        made += 1

    # A deck with only picture-less content layouts still needs ≥1 slide;
    # drop nothing else — save and report.
    prs.save(out)
    return {'path': out, 'slideCount': made, 'kind': 'pptx'}


# ── Charts ────────────────────────────────────────────────────────────────


_CHART_KINDS = ('line', 'bar', 'pie', 'scatter', 'hist')


def render_chart(
    path: str,
    kind: str,
    series: Any,
    labels: Any = None,
    title: str = '',
    xlabel: str = '',
    ylabel: str = '',
    workspace: str = '',
) -> dict[str, Any]:
    """Render line/bar/pie/scatter/hist chart(s) to PNG.

    ``series``: [[num, ...], ...] — one inner list per plotted set.
    ``labels``: category/x labels (bar/pie) — len must match the data.
    """
    out = _bind_write(path, workspace, ('.png', '.jpg', '.jpeg', '.webp'))
    k = (kind or '').strip().lower()
    if k not in _CHART_KINDS:
        raise ValueError(f'kind must be one of {", ".join(_CHART_KINDS)}.')
    data = _as_list(series)
    if not data or not all(isinstance(row, (list, tuple)) for row in data):
        # Allow a flat numeric list as a single series.
        if data and all(isinstance(v, (int, float)) for v in data):
            data = [data]
        else:
            raise ValueError('series must be a list of numeric lists.')
    rows = [[float(v) for v in row] for row in data]

    plt = _ensure_matplotlib()
    fig, ax = plt.subplots(figsize=(8, 5), dpi=150)
    try:
        lab = [str(x) for x in _as_list(labels)] if labels else None
        if k == 'line':
            for r in rows:
                ax.plot(r, marker='o', linewidth=2)
        elif k == 'bar':
            xs = range(len(rows[0]))
            width = 0.8 / max(1, len(rows))
            for idx, r in enumerate(rows):
                ax.bar([x + idx * width for x in xs], r, width=width)
            if lab:
                ax.set_xticks([x + 0.4 - width / 2 for x in xs])
                ax.set_xticklabels(lab)
        elif k == 'pie':
            if len(rows) != 1:
                raise ValueError('pie takes exactly one series.')
            if lab and len(lab) != len(rows[0]):
                raise ValueError('labels length must match the pie series.')
            ax.pie(rows[0], labels=lab or None, autopct='%1.1f%%', startangle=90)
            ax.axis('equal')
        elif k == 'scatter':
            if len(rows) != 2:
                raise ValueError('scatter takes exactly two series: x then y.')
            if len(rows[0]) != len(rows[1]):
                raise ValueError('scatter x/y lengths differ.')
            ax.scatter(rows[0], rows[1], s=28)
        else:  # hist
            for r in rows:
                ax.hist(r, bins=min(24, max(6, len(r) // 2 or 6)), alpha=0.7)
        if title:
            ax.set_title(title)
        if xlabel and k != 'pie':
            ax.set_xlabel(xlabel)
        if ylabel and k != 'pie':
            ax.set_ylabel(ylabel)
        if k in ('line', 'bar', 'hist'):
            ax.grid(True, alpha=0.25)
        fig.tight_layout()
        fig.savefig(out)
    finally:
        plt.close(fig)
    return {'path': out, 'kind': k, 'seriesCount': len(rows)}


# ── Video ─────────────────────────────────────────────────────────────────


def render_video(
    path: str,
    frames: Any = None,
    fps: int = 12,
    hold_last_ms: int = 400,
    workspace: str = '',
) -> dict[str, Any]:
    """Assemble an MP4 from image paths (one per frame).

    ``frames``: list of image file paths (workspace-relative or absolute).
    Every frame gets equal duration (1/fps); the final frame is held
    ``hold_last_ms`` so the video doesn't cut off abruptly.
    """
    out = _bind_write(path, workspace, ('.mp4',))
    items = _as_list(frames)
    if not items:
        raise ValueError('frames must list at least one image path.')


    import numpy as np
    from PIL import Image

    arrays: list[Any] = []
    size: tuple[int, int] | None = None
    for raw in items:
        p = str(raw)
        bound, err = bind_path(p, workspace, for_write=False)
        if err or bound is None or not bound.exists():
            raise ValueError(f'Frame not found / outside workspace: {p}')
        img = Image.open(bound).convert('RGB')
        if size is None:
            # Even dimensions keep encoders happy.
            size = (img.width - img.width % 2, img.height - img.height % 2)
        if (img.width, img.height) != size:
            img = img.resize(size)
        arrays.append(np.asarray(img, dtype=np.uint8))

    import imageio.v2 as iio

    writer = iio.get_writer(out, fps=int(fps), codec='libx264', quality=8, macro_block_size=None)
    try:
        for idx, frame in enumerate(arrays):
            writer.append_data(frame)
            if idx == len(arrays) - 1 and hold_last_ms > 0:
                extra = max(1, round(int(hold_last_ms) / 1000 * int(fps)))
                for _ in range(extra - 1):
                    writer.append_data(frame)
    finally:
        writer.close()

    count = len(arrays)
    dur = count / max(1, int(fps)) + (int(hold_last_ms) / 1000 if count else 0)
    return {
        'path': out,
        'frameCount': count,
        'fps': int(fps),
        'durationSec': round(dur, 2),
        'size': list(size or ()),
    }


# ── Circuits ──────────────────────────────────────────────────────────────

# schemdraw element table: JSON type → schemdraw.elements class name.
_CIRCUIT_ELEMENTS = {
    'resistor': 'Resistor',
    'r': 'Resistor',
    'capacitor': 'Capacitor',
    'c': 'Capacitor',
    'inductor': 'Inductor',
    'l': 'Inductor',
    'diode': 'Diode',
    'd': 'Diode',
    'led': 'LED',
    'zener': 'Zener',
    'battery': 'Battery',
    'battery1': 'Battery',
    'voltage': 'SourceV',
    'sourcev': 'SourceV',
    'v': 'SourceV',
    'current': 'SourceI',
    'sourcei': 'SourceI',
    'i': 'SourceI',
    'switch': 'Switch',
    'sw': 'SwitchSpst',
    'lamp': 'Lamp',
    'fuse': 'Fuse',
    'speaker': 'Speaker',
    'motor': 'Motor',
    'opamp': 'Opamp',
    'ground': 'Ground',
    'gnd': 'Ground',
    'line': 'Line',
    'wire': 'Line',
}


def draw_circuit(path: str, elements: Any, title: str = '', workspace: str = '') -> dict[str, Any]:
    """Draw a daisy-chain schematic to PNG with schemdraw.

    ``elements``: [{"type": "battery"|"resistor"|..., "label": "10V",
                    "dir": "right|left|up|down"}, ...] drawn left-to-right
    in order; "ground" anchors the chain. This covers the classic
    series-loop diagrams (PSUs, dividers, LED drivers) the harness is
    asked to explain; heavier topologies stay future work.
    """
    out = _bind_write(path, workspace, ('.png',))
    specs = _as_list(elements)
    if not specs:
        raise ValueError('elements must describe at least one component.')

    import matplotlib

    matplotlib.use('Agg', force=True)
    import schemdraw
    import schemdraw.elements as elm

    d = schemdraw.Drawing(unit=2.5)
    placed = 0
    for i, raw in enumerate(specs):
        if isinstance(raw, str):
            raw = {'type': raw}
        if not isinstance(raw, dict):
            raise ValueError(f'elements[{i}] must be an object.')
        etype = str(raw.get('type', '')).strip().lower()
        cls_name = _CIRCUIT_ELEMENTS.get(etype)
        if not cls_name:
            known = ', '.join(sorted(set(_CIRCUIT_ELEMENTS)))
            raise ValueError(f'Unknown element type "{etype}". Known: {known}')
        elm_cls = getattr(elm, cls_name, None)
        if elm_cls is None:
            raise ValueError(f'schemdraw has no element "{cls_name}".')
        part = elm_cls()
        direction = str(raw.get('dir', '')).strip().lower()
        if direction in ('up', 'down', 'left', 'right'):
            part = getattr(part, direction)()
        label = str(raw.get('label', '') or '').strip()
        if label:
            part = part.label(label)
        d += part
        placed += 1

    if title:
        d += elm.Label().label(title).at((0, -2))
    d.save(out, dpi=150)
    return {'path': out, 'elementCount': placed}


# ── Interactive HTML artifacts ────────────────────────────────────────────


def create_html_artifact(
    path: str,
    html: str,
    title: str = '',
    workspace: str = '',
) -> dict[str, Any]:
    """Write a SELF-CONTAINED interactive HTML explainer/animation.

    The model authors the full document (inline CSS + JS, canvas/SVG
    animations, step-through explainers — anything a browser can run
    offline). August renders it live in the right-side panel viewer.
    No network resources: everything inline so the artifact works forever.
    """
    out = _bind_write(path, workspace, ('.html', '.htm'))
    if not isinstance(html, str) or not html.strip():
        raise ValueError('html must be a non-empty document string.')

    # Guard against accidental external dependence breaking the offline
    # viewer: warn (not fail) when remote resources are referenced.
    lowered = html.lower()
    external = [
        token
        for token in ('http://', 'https://', '//cdn.', 'src="http')
        if token in lowered
    ]

    doc = html
    if title and '<title' not in lowered:
        safe_title = title.replace('<', '&lt;').replace('>', '&gt;')
        doc = doc.replace('<head>', f'<head>\n<title>{safe_title}</title>', 1) \
            if '<head>' in lowered else f'<title>{safe_title}</title>\n{doc}'

    bound = Path(out)
    bound.write_text(doc, encoding='utf-8')
    return {
        'path': out,
        'bytes': len(doc.encode('utf-8')),
        'externalRefs': external,
    }
